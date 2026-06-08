import streamlit as st
import os
import sqlite3
import chardet
import re
import olefile
from pathlib import Path

try:
    from tkinter import Tk, filedialog
except:
    Tk = None
    filedialog = None

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# ======================全局配置======================
DB_NAME = "knowledge.db"
TEMP_UPLOAD_DIR = "upload_files"
PAGE_SIZE = 50
FILE_PAGE = 30
PDF_MAX_CONTENT = 3000

SUFFIX_LIST = (".pdf", ".doc", ".docx", ".ppt", ".pptx", ".txt", ".md", ".html", ".htm")

# ======================数据库初始化======================
def init_database():
    conn = sqlite3.connect(DB_NAME)
    cur = conn.cursor()
    cur.execute('''CREATE TABLE IF NOT EXISTS knowledge_data (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        file_name TEXT,
        file_type TEXT,
        page_num INTEGER,
        content TEXT,
        full_path TEXT
    )''')
    try:
        cur.execute("ALTER TABLE knowledge_data ADD COLUMN full_path TEXT;")
    except Exception:
        pass
    cur.execute('''CREATE VIRTUAL TABLE IF NOT EXISTS knowledge_index 
    USING fts5(file_name, file_type, page_num, content)''')
    conn.commit()
    conn.close()

def get_db_exist_filename():
    conn = sqlite3.connect(DB_NAME)
    res = conn.execute("SELECT DISTINCT file_name FROM knowledge_data").fetchall()
    conn.close()
    return {row[0] for row in res}

# ======================解析函数======================
def parse_html(file_path):
    fname = os.path.basename(file_path)
    try:
        with open(file_path, "rb") as f:
            raw = f.read()
        enc = chardet.detect(raw)["encoding"] or "utf-8"
        html = raw.decode(enc, errors="ignore")
    except:
        html = ""
    clean = re.sub(r"<[^>]+>", " ", html)
    clean = re.sub(r"\s+", " ", clean).strip()
    res_list = []
    cut_len = 800
    idx = 1
    while len(clean) > 0:
        seg = clean[:cut_len].strip()
        if seg:
            res_list.append({
                "file_name": fname, "file_type": "HTML", "page_num": idx,
                "content": seg, "full_path": file_path
            })
        clean = clean[cut_len:]
        idx += 1
    return res_list

def parse_text_file(file_path):
    fname = os.path.basename(file_path)
    with open(file_path, "rb") as f:
        raw = f.read()
    code_order = ["gb18030", "gbk", "gb2312", "utf-8", "cp1252", "latin-1"]
    txt = ""
    for cd in code_order:
        try:
            txt = raw.decode(cd)
            break
        except Exception:
            continue
    if not txt:
        enc_auto = chardet.detect(raw)["encoding"] or "utf-8"
        txt = raw.decode(enc_auto, errors="replace")
    txt = txt.replace("\ufeff", "")
    lines = [i.strip() for i in txt.splitlines() if i.strip()]
    res = []
    for num, line in enumerate(lines, start=1):
        res.append({
            "file_name": fname, "file_type": "TXT", "page_num": num,
            "content": line, "full_path": file_path
        })
    return res

def parse_pdf(file_path):
    fname = os.path.basename(file_path)
    res = []
    try:
        reader = PdfReader(file_path)
        total_p = len(reader.pages)
        for pg_idx in range(total_p):
            page = reader.pages[pg_idx]
            raw_txt = page.extract_text() or ""
            raw_txt = raw_txt.encode("utf8", errors="replace").decode("utf8")
            cont = raw_txt.strip()
            if not cont:
                continue
            if len(cont) > PDF_MAX_CONTENT:
                cont = cont[:PDF_MAX_CONTENT]
            res.append({
                "file_name": fname, "file_type": "PDF", "page_num": pg_idx + 1,
                "content": cont, "full_path": file_path
            })
    except Exception as e:
        st.warning(f"PDF解析异常 {fname}:{str(e)}")
    return res

def parse_docx(file_path):
    fname = os.path.basename(file_path)
    res = []
    try:
        doc = Document(file_path)
        full_text = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        all_str = "\n".join(full_text)
        split_size = 1500
        piece_list = [all_str[i:i+split_size] for i in range(0, len(all_str), split_size)]
        for no, piece in enumerate(piece_list, start=1):
            res.append({
                "file_name": fname, "file_type": "WORD", "page_num": no,
                "content": piece, "full_path": file_path
            })
    except Exception:
        pass
    return res

def parse_old_doc(file_path):
    fname = os.path.basename(file_path)
    res = []
    try:
        ole = olefile.OleFileIO(file_path)
        stream = ole.openstream('WordDocument')
        raw_data = stream.read()
        ole.close()
        txt = ""
        try:
            txt = raw_data.decode("gb18030", errors="replace")
        except:
            txt = raw_data.decode("utf-8", errors="replace")
        txt = txt.replace("\x00", "")
        paragraphs = [p.strip() for p in txt.splitlines() if p.strip()]
        all_text = "".join(paragraphs)
        split_size = 1500
        piece_list = [all_text[i:i+split_size] for i in range(0, len(all_text), split_size)]
        for no, piece in enumerate(piece_list, start=1):
            res.append({
                "file_name": fname, "file_type": "OLD_DOC", "page_num": no,
                "content": piece, "full_path": file_path
            })
    except Exception as e:
        st.warning(f"老DOC解析失败 {fname}:{str(e)}")
    return res

def parse_pptx(file_path):
    fname = os.path.basename(file_path)
    res = []
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shp in slide.shapes:
                if hasattr(shp, "text") and shp.text.strip():
                    full_text.append(shp.text.strip())
        all_str = "\n".join(full_text)
        split_size = 1500
        piece_list = [all_str[i:i+split_size] for i in range(0, len(all_str), split_size)]
        for no, piece in enumerate(piece_list, start=1):
            res.append({
                "file_name": fname, "file_type": "PPT", "page_num": no,
                "content": piece, "full_path": file_path
            })
    except Exception:
        pass
    return res

# ======================批量入库======================
def batch_save_files(file_list, bar):
    exist_set = get_db_exist_filename()
    conn = sqlite3.connect(DB_NAME)
    cur = conn.cursor()
    total_cnt = len(file_list)
    succ = 0
    skip = 0
    for idx, fp in enumerate(file_list):
        fn = os.path.basename(fp)
        if fn in exist_set:
            skip += 1
            bar.progress((idx+1)/total_cnt, text=f"跳过重复:{fn}")
            continue
        ext = os.path.splitext(fp.lower())[-1]
        parse_result = []
        if ext == ".pdf":
            parse_result = parse_pdf(fp)
        elif ext == ".docx":
            parse_result = parse_docx(fp)
        elif ext == ".doc":
            parse_result = parse_old_doc(fp)
        elif ext in (".ppt", ".pptx"):
            parse_result = parse_pptx(fp)
        elif ext in (".txt", ".md"):
            parse_result = parse_text_file(fp)
        elif ext in (".html", ".htm"):
            parse_result = parse_html(fp)

        for item in parse_result:
            cur.execute('INSERT INTO knowledge_data(file_name,file_type,page_num,content,full_path) VALUES (?,?,?,?,?)',
                        (item["file_name"], item["file_type"], item["page_num"], item["content"], item["full_path"]))
            cur.execute('INSERT INTO knowledge_index(file_name,file_type,page_num,content) VALUES (?,?,?,?)',
                        (item["file_name"], item["file_type"], item["page_num"], item["content"]))
        exist_set.add(fn)
        succ += 1
        bar.progress((idx+1)/total_cnt, text=f"成功:{succ}")
    conn.commit()
    conn.close()
    return succ, skip, total_cnt

# ======================数据查询======================
def get_all_file_total():
    conn = sqlite3.connect(DB_NAME)
    data = conn.execute("SELECT DISTINCT file_name,full_path FROM knowledge_data ORDER BY file_name").fetchall()
    conn.close()
    return data

def get_file_page(all_list, page):
    total = len(all_list)
    start = (page-1)*FILE_PAGE
    end = start + FILE_PAGE
    page_data = all_list[start:end]
    max_page = (total + FILE_PAGE -1)//FILE_PAGE if total>0 else 1
    return page_data, max_page, total

def get_item_page_data(page):
    off = (page-1)*PAGE_SIZE
    conn = sqlite3.connect(DB_NAME)
    rows = conn.execute('SELECT file_name,file_type,page_num,content FROM knowledge_data ORDER BY file_name,page_num LIMIT ? OFFSET ?',
                        (PAGE_SIZE, off)).fetchall()
    total_all = conn.execute("SELECT COUNT(*) FROM knowledge_data").fetchone()[0]
    conn.close()
    max_page = (total_all + PAGE_SIZE -1)//PAGE_SIZE if total_all>0 else 1
    return rows, total_all, max_page

def open_full_file(fp):
    try:
        ext = Path(fp).suffix.lower()
        if ext in (".txt",".md",".html",".htm"):
            with open(fp,"rb") as f:
                raw = f.read()
            cd_list = ["gb18030","gbk","gb2312","utf-8","cp1252"]
            for c in cd_list:
                try:
                    return raw.decode(c)
                except Exception:
                    continue
            return "编码解析失败"
        elif ext == ".pdf":
            rd = PdfReader(fp)
            full = ""
            for p in rd.pages:
                t = p.extract_text() or ""
                full += t + "\n\n"
            return full
        elif ext in (".doc",".docx"):
            doc = Document(fp)
            return "\n".join([p.text for p in doc.paragraphs])
        elif ext in (".ppt",".pptx"):
            ppt = Presentation(fp)
            arr = []
            for s in ppt.slides:
                for sh in s.shapes:
                    if hasattr(sh,"text"):
                        arr.append(sh.text)
            return "\n".join(arr)
        else:
            return "暂不支持预览"
    except Exception as e:
        return f"打开失败:{str(e)}"

def delete_one_file(filename):
    conn = sqlite3.connect(DB_NAME)
    cur = conn.cursor()
    cur.execute("DELETE FROM knowledge_data WHERE file_name=?", (filename,))
    cur.execute("DELETE FROM knowledge_index WHERE file_name=?", (filename,))
    conn.commit()
    conn.close()

def search_content(key):
    key = re.sub(r'[\\/*?:"<>|]', "", key.strip())
    if not key:
        return []
    try:
        conn = sqlite3.connect(DB_NAME)
        res = conn.execute("SELECT file_name,file_type,page_num,content FROM knowledge_index WHERE content MATCH ?", (key,)).fetchall()
        conn.close()
        return res
    except Exception:
        return []

def pick_folder():
    if not Tk:
        return ""
    try:
        root = Tk()
        root.withdraw()
        root.attributes("-topmost", True)
        path = filedialog.askdirectory()
        root.destroy()
        return path
    except Exception:
        return ""

def scan_folder_file(folder_path):
    if not os.path.isdir(folder_path):
        return []
    res = []
    for root, _, files in os.walk(folder_path):
        for name in files:
            if name.lower().endswith(SUFFIX_LIST):
                res.append(os.path.join(root, name))
    return res

# ======================导出======================
def export_md(search_res):
    md_text = "# 检索结果\n"
    for idx, item in enumerate(search_res, 1):
        md_text += f"{idx}. {item[0]} 段落{item[2]}\n{item[3]}\n---\n"
    return md_text

def export_txt(search_res):
    t = "检索结果\n====================\n"
    for idx, item in enumerate(search_res, 1):
        t += f"{idx}｜{item[0]} P{item[2]}\n{item[3]}\n\n"
    return t

def export_word(search_res):
    doc = Document()
    doc.add_heading("检索结果", 0)
    for it in search_res:
        doc.add_paragraph(f"文件名：{it[0]} 分段：{it[2]}")
        doc.add_paragraph(it[3])
    doc.save("temp_out.docx")
    with open("temp_out.docx", "rb") as f:
        b = f.read()
    os.remove("temp_out.docx")
    return b

# ======================主界面======================
def main():
    st.set_page_config(page_title="⚡极速知识库", layout="wide")
    st.title("⚡ 极速知识库")
    init_database()
    os.makedirs(TEMP_UPLOAD_DIR, exist_ok=True)

    if "curr_page" not in st.session_state:
        st.session_state.curr_page = 1
    if "file_page" not in st.session_state:
        st.session_state.file_page = 1
    if "local_folder" not in st.session_state:
        st.session_state.local_folder = ""
    if "view_type" not in st.session_state:
        st.session_state.view_type = "文件模式"
    if "open_full_file_info" not in st.session_state:
        st.session_state.open_full_file_info = None

    if st.session_state.open_full_file_info is not None:
        fname, fpath = st.session_state.open_full_file_info
        st.subheader(f"📄完整原文：{fname}")
        full_txt = open_full_file(fpath)
        st.text_area("文件全部内容", full_txt, height=550)
        if st.button("关闭预览返回列表"):
            st.session_state.open_full_file_info = None
            st.rerun()
        return

    # 顶部常驻检索
    with st.expander("🔍 全文检索", expanded=True):
        search_key = st.text_input("输入关键词")
        if search_key:
            res_data = search_content(search_key)
            if res_data:
                st.success(f"命中 {len(res_data)} 条结果")
                md_str = export_md(res_data)
                txt_str = export_txt(res_data)
                word_byte = export_word(res_data)
                d1, d2, d3 = st.columns(3)
                with d1:
                    st.download_button("导出MD", md_str, "检索结果.md")
                with d2:
                    st.download_button("导出TXT", txt_str, "检索结果.txt")
                with d3:
                    st.download_button("导出Word", word_byte, "检索结果.docx")
                for item in res_data:
                    with st.expander(f"{item[0]}｜分段{item[2]}"):
                        st.write(item[3])
            else:
                st.info("无匹配内容")
    st.divider()

    # 侧边栏
    with st.sidebar:
        st.header("📂文件管理")
        st.subheader("本地文件上传")
        upload_list = st.file_uploader("支持pdf/doc/docx/ppt/pptx/txt/md/html/htm",
                                      accept_multiple_files=True, type=[i.strip(".") for i in SUFFIX_LIST])
        if upload_list:
            temp_save = []
            for upf in upload_list:
                save_p = os.path.join(TEMP_UPLOAD_DIR, upf.name)
                with open(save_p, "wb") as wf:
                    wf.write(upf.getbuffer())
                temp_save.append(save_p)
            if st.button("✅确认入库"):
                bar = st.progress(0)
                succ, skip, _ = batch_save_files(temp_save, bar)
                bar.empty()
                st.success(f"新增{succ}个文件，跳过重复{skip}个")

        st.divider()
        st.subheader("文件夹批量导入")
        input_dir = st.text_input("目录路径", value=st.session_state.local_folder)
        if st.button("选择目录"):
            sel_p = pick_folder()
            if sel_p:
                st.session_state.local_folder = sel_p
                st.rerun()
        if input_dir and st.button("开始批量导入"):
            all_f = scan_folder_file(input_dir)
            if all_f:
                bar = st.progress(0)
                succ, skip, _ = batch_save_files(all_f, bar)
                bar.empty()
                st.success(f"共扫描{len(all_f)}，成功{succ}，重复跳过{skip}")
            else:
                st.warning("目录无匹配格式文件")

        st.divider()
        st.subheader("删除库内文件")
        all_file_data = get_all_file_total()
        name_arr = [i[0] for i in all_file_data]
        if name_arr:
            del_sel = st.selectbox("选择删除文件", name_arr)
            if st.button("确认删除"):
                delete_one_file(del_sel)
                st.success("删除成功")
                st.rerun()

    # 浏览模式切换
    sel_view = st.radio("浏览方式：", ["文件模式", "条目分页模式"], horizontal=True)
    st.session_state.view_type = sel_view
    st.divider()

    # 文件浏览（分页）
    if st.session_state.view_type == "文件模式":
        with st.expander("📁全部文件列表(点击打开完整原文)", expanded=True):
            all_file_raw = get_all_file_total()
            file_page_list, max_fpage, total_all = get_file_page(all_file_raw, st.session_state.file_page)
            st.info(f"库总文件：{total_all}｜当前页 {st.session_state.file_page}/{max_fpage}")

            c1, c2, c3, c4, c5 = st.columns([1,1,1.8,1,1])
            with c1:
                if st.button("首页"):
                    st.session_state.file_page = 1
                    st.rerun()
            with c2:
                if st.button("上一页"):
                    st.session_state.file_page = max(1, st.session_state.file_page - 1)
                    st.rerun()
            with c3:
                jump_f = st.number_input("跳转", value=st.session_state.file_page, min_value=1, max_value=max_fpage)
                if jump_f != st.session_state.file_page:
                    st.session_state.file_page = jump_f
                    st.rerun()
            with c4:
                if st.button("下一页"):
                    st.session_state.file_page = min(max_fpage, st.session_state.file_page + 1)
                    st.rerun()
            with c5:
                if st.button("末页"):
                    st.session_state.file_page = max_fpage
                    st.rerun()

            for fn, fp in file_page_list:
                cc1, cc2 = st.columns([5,1])
                with cc1:
                    st.markdown(f"**{fn}**")
                with cc2:
                    if st.button("打开全文", key="f_"+fn):
                        st.session_state.open_full_file_info = (fn, fp)
                        st.rerun()
                st.divider()

    # 条目浏览
    else:
        with st.expander("📄分段条目浏览", expanded=True):
            rows, total_item, max_p = get_item_page_data(st.session_state.curr_page)
            st.info(f"全库总条目：{total_item}｜当前页 {st.session_state.curr_page}/{max_p}")

            b1, b2, b3, b4, b5 = st.columns([1,1,1.8,1,1])
            with b1:
                if st.button("⏮️首页"):
                    st.session_state.curr_page = 1
                    st.rerun()
            with b2:
                if st.button("⬅️上一页"):
                    st.session_state.curr_page = max(1, st.session_state.curr_page - 1)
                    st.rerun()
            with b3:
                jump_p = st.number_input("页码", value=st.session_state.curr_page, min_value=1, max_value=max_p)
                if jump_p != st.session_state.curr_page:
                    st.session_state.curr_page = jump_p
                    st.rerun()
            with b4:
                if st.button("➡️下一页"):
                    st.session_state.curr_page = min(max_p, st.session_state.curr_page + 1)
                    st.rerun()
            with b5:
                if st.button("⏭️末页"):
                    st.session_state.curr_page = max_p
                    st.rerun()

            for fname, ftype, pno, cont in rows:
                with st.expander(f"{fname}｜{ftype}｜分段{pno}"):
                    st.write(cont)
                st.divider()

if __name__ == "__main__":
    main()
